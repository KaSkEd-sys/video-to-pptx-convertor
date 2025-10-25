import os
import cv2
from pptx import Presentation
from pptx.util import Inches
import sys


def print_progress_bar(iteration, total, prefix='', suffix='', length=50, fill='█'):
    """Displays a progress bar in the console."""
    percent = f"{100 * (iteration / float(total)):.1f}"
    filled_length = int(length * iteration // total)
    bar = fill * filled_length + '-' * (length - filled_length)
    sys.stdout.write(f'\r{prefix} |{bar}| {percent}% {suffix}')
    sys.stdout.flush()
    if iteration == total:
        print()


def extract_frames(video_path="video.mp4", output_folder="base", frame_step=10):
    """Extracts frames from a video and saves them into a folder."""
    if not os.path.exists(video_path):
        print(f"❌ Video '{video_path}' not found.")
        return 0

    os.makedirs(output_folder, exist_ok=True)

    cap = cv2.VideoCapture(video_path)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    count = 0
    saved = 0

    print(f"🎞 Extracting frames from {video_path}...")

    while True:
        ret, frame = cap.read()
        if not ret:
            break
        if count % frame_step == 0:  # save every Nth frame
            filename = os.path.join(output_folder, f"frame_{saved:05d}.jpg")
            cv2.imwrite(filename, frame)
            saved += 1
            print_progress_bar(count, total_frames, prefix='Extracting frames', suffix='Done', length=40)
        count += 1

    cap.release()
    print(f"\n✅ {saved} frames extracted to folder '{output_folder}'.")
    return saved


def images_to_ppt(frames_folder="base", output_ppt="frames_to_slides.pptx"):
    """Creates a PowerPoint presentation from images."""
    if not os.path.exists(frames_folder):
        print(f"❌ Folder '{frames_folder}' not found.")
        return

    images = [f for f in os.listdir(frames_folder) if f.lower().endswith((".jpg", ".png"))]
    images.sort()

    if not images:
        print(f"⚠️ No images found in folder '{frames_folder}'.")
        return

    print(f"📸 Found {len(images)} frames. Creating presentation...")

    prs = Presentation()

    for i, img in enumerate(images, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        img_path = os.path.join(frames_folder, img)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)
        print_progress_bar(i, len(images), prefix='Adding slides', suffix='Done', length=40)

    prs.save(output_ppt)
    print(f"\n✅ Presentation successfully created: {os.path.abspath(output_ppt)}")


if __name__ == "__main__":
    video_path = "video.mp4"
    frames_folder = "base"
    output_ppt = "frames_to_slides.pptx"

    # 1️⃣ Extract frames from video
    extracted = extract_frames(video_path, frames_folder, frame_step=10)

    # 2️⃣ If frames were extracted, create the PowerPoint presentation
    if extracted > 0:
        images_to_ppt(frames_folder, output_ppt)
